#!/usr/bin/env python3.10
"""
热点刀锋选题报告PPT生成器 (v1.1)
===================================
用法: python3.10 scripts/generate_topic_report_pptx.py --data /tmp/hotlist_data.json --history ~/.hermes/skills/productivity/hotspot-blade/hotspot-blade-history.json --output /mnt/c/Users/yingm/OneDrive/Desktop/选题报告_{date}.pptx

需要: python-pptx (pip install python-pptx)
必须用python3.10执行（如有openpyxl依赖冲突）

输入: hotlist_scraper.py输出的JSON文件（5平台热榜数据）
     topics JSON文件（预评分话题数组，每话题含title/score/grade/scores/source/heat/angle/object_detail/split_detail）
输出: 8页PPT选题报告（深色主题、评分条、推荐排序）

v1.1更新 (2026-06-14):
  - 推荐排序页(slide 7)接受topics中的'reason'字段，不再使用硬编码文本
  - 支持自定义reason字段（每话题一条），为空时自动根据角度生成

陷阱:
  - 字符串中的中文引号(如"东西")会导致SyntaxError，必须用交替引号或转义
  - PPTX保存路径必须为Windows桌面路径 (/mnt/c/Users/...)
"""

import json
import sys
import argparse
from datetime import datetime
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
except ImportError:
    print("Error: python-pptx not installed. Run: pip install python-pptx")
    sys.exit(1)


# ── Color palette ──
BG_DARK = RGBColor(0x1a, 0x1a, 0x2e)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
ACCENT_BLUE = RGBColor(0x4e, 0xca, 0xf8)
ACCENT_GREEN = RGBColor(0x4c, 0xaf, 0x50)
ACCENT_YELLOW = RGBColor(0xFF, 0xC1, 0x07)
ACCENT_RED = RGBColor(0xF4, 0x43, 0x36)
ACCENT_ORANGE = RGBColor(0xFF, 0x98, 0x00)
GRAY = RGBColor(0xBB, 0xBB, 0xBB)


def set_bg(slide, color=BG_DARK):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, left, top, width, height, text, font_size=18,
                bold=False, color=WHITE, alignment=PP_ALIGN.LEFT,
                font_name='Microsoft YaHei'):
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return tf


def add_paragraph(text_frame, text, font_size=18, bold=False, color=WHITE,
                  alignment=PP_ALIGN.LEFT, space_before=Pt(6),
                  font_name='Microsoft YaHei'):
    p = text_frame.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    if space_before:
        p.space_before = space_before
    return p


def add_rich_textbox(slide, left, top, width, height):
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    return tf


def add_paragraph_rich(tf, runs_list, alignment=PP_ALIGN.LEFT,
                       space_before=Pt(4)):
    """runs_list: list of (text, font_size, bold, color) tuples"""
    p = tf.add_paragraph()
    p.alignment = alignment
    if space_before:
        p.space_before = space_before
    for text, font_size, bold, color in runs_list:
        run = p.add_run()
        run.text = text
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.color.rgb = color
        run.font.name = 'Microsoft YaHei'
    return p


def draw_score_bar(tf, label, score, max_score=10, bar_width_inches=6):
    """Append a score bar line to a text frame"""
    filled = int(score / max_score * 10)
    empty = 10 - filled
    bar = chr(0x2588) * filled + chr(0x2591) * empty
    bar_color = ACCENT_GREEN if score >= 8 else (ACCENT_YELLOW if score >= 6 else ACCENT_RED)
    p = tf.add_paragraph()
    p.space_before = Pt(3)
    run1 = p.add_run()
    run1.text = label + '  '
    run1.font.size = Pt(14)
    run1.font.bold = True
    run1.font.color.rgb = WHITE
    run1.font.name = 'Microsoft YaHei'
    run2 = p.add_run()
    run2.text = f'{score}/10  '
    run2.font.size = Pt(14)
    run2.font.bold = True
    run2.font.color.rgb = bar_color
    run2.font.name = 'Microsoft YaHei'
    run3 = p.add_run()
    run3.text = bar
    run3.font.size = Pt(10)
    run3.font.color.rgb = bar_color
    run3.font.name = 'Consolas'


def draw_topic_block(slide, topic, left, top):
    """Draw one topic scoring block on a slide"""
    grade_color = ACCENT_GREEN if topic['grade'] == 'S' else ACCENT_YELLOW
    tf = add_rich_textbox(slide, left, top, 5.5, 0.5)
    p0 = tf.paragraphs[0]
    p0.text = ''
    add_paragraph_rich(tf, [
        (topic['title'], 18, True, WHITE),
        (f'    总分: {topic["score"]}/10  {topic["grade"]}级', 16, True, grade_color),
    ])
    tf2 = add_rich_textbox(slide, left, top + 0.6, 5.5, 2.0)
    p0 = tf2.paragraphs[0]
    p0.text = ''
    for label, score in topic['scores'].items():
        draw_score_bar(tf2, label, score)
    tf3 = add_rich_textbox(slide, left, top + 2.7, 5.5, 1.5)
    p0 = tf3.paragraphs[0]
    p0.text = ''
    add_paragraph_rich(tf3, [
        ('来源: ', 14, True, GRAY), (topic['source'], 14, False, WHITE),
        ('  |  热度: ', 14, True, GRAY), (topic['heat'], 14, False, ACCENT_YELLOW),
    ])
    add_paragraph_rich(tf3, [
        ('角度: ', 14, True, GRAY), (topic['angle'], 14, False, WHITE),
    ])
    add_paragraph_rich(tf3, [
        ('物件: ', 14, True, GRAY), (topic.get('object_detail', ''), 14, False, ACCENT_ORANGE),
    ])


def build_report(topics, today=None, platform_counts=None):
    """Build the 8-slide PPT report.

    Args:
        topics: list of dict with fields {title, score, grade, scores:{},
                source, heat, angle, object_detail, split_detail}
                Optional: 'reason' field for slide 7 ranking descriptions
        today: date string like '2026-06-14', defaults to today
        platform_counts: dict like {'zhihu': 30, 'weibo': 50, ...}

    Returns:
        pptx.Presentation object
    """
    if today is None:
        today = datetime.now().strftime('%Y-%m-%d')

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # ── Slide 1: Cover ──
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_textbox(slide, 1, 2, 11, 1.5, '热点刀锋选题报告', font_size=44, bold=True, color=ACCENT_BLUE, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, 1, 3.5, 11, 0.6, today, font_size=24, color=GRAY, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, 1, 4.3, 11, 0.5, '5平台热榜  ->  5维度评分  ->  精选5个S/A级话题', font_size=18, color=GRAY, alignment=PP_ALIGN.CENTER)

    # ── Slide 2: Data Overview ──
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_textbox(slide, 0.8, 0.3, 11, 0.6, '数据总览', font_size=32, bold=True, color=ACCENT_BLUE)
    tf = add_rich_textbox(slide, 0.8, 1.2, 11, 5)
    p0 = tf.paragraphs[0]
    p0.text = ''
    now_str = datetime.now().strftime('%Y-%m-%d %H:%M')
    add_paragraph_rich(tf, [('抓取时间: ', 18, True, GRAY), (now_str, 18, False, WHITE)])
    if platform_counts:
        total = sum(platform_counts.values())
        add_paragraph_rich(tf, [('数据来源: ', 18, True, GRAY), ('知乎热榜  ·  微博热搜  ·  B站热门  ·  36氪热榜  ·  百度热搜', 18, False, WHITE)])
        add_paragraph_rich(tf, [('总抓取量: ', 18, True, GRAY), (f'{total}条原始数据', 18, False, ACCENT_GREEN)], space_before=Pt(8))
    add_paragraph_rich(tf, [('排除: ', 18, True, GRAY), ('近7天已写话题', 18, False, WHITE)], space_before=Pt(4))
    add_paragraph_rich(tf, [('评分体系: ', 18, True, GRAY), ('钱包距离35% + 反驳成本25% + 物件锚点15% + 头条适配15% + 天然分裂10%', 18, False, WHITE)], space_before=Pt(4))
    add_paragraph_rich(tf, [('交叉验证加分: ', 18, True, GRAY), ('同一话题在2+平台同时出现  额外+0.5分', 18, False, ACCENT_YELLOW)], space_before=Pt(4))

    # ── Slide 3: Scoring System ──
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_textbox(slide, 0.8, 0.3, 11, 0.6, '评分体系说明', font_size=32, bold=True, color=ACCENT_BLUE)
    tf = add_rich_textbox(slide, 0.8, 1.1, 11, 5.5)
    p0 = tf.paragraphs[0]
    p0.text = ''
    dims = [
        ('钱包距离', ACCENT_GREEN, 35, '这个话题跟读者的钱包有多近？',
         '9-10: 直接关乎工资/房价/物价  |  5-6: 间接影响  |  1-2: 跟钱包无关'),
        ('反驳成本', ACCENT_YELLOW, 25, '越低越好！读者能不能一句话反驳？',
         '9-10: 一句话反驳  |  5-6: 需要举例子  |  1-2: 无法反驳'),
        ('物件锚点', ACCENT_ORANGE, 15, '有没有读者生活中见过的具体"东西"？',
         '9-10: 天天用/见  |  5-6: 知道但不接触  |  1-2: 纯抽象概念'),
        ('头条适配性', ACCENT_BLUE, 15, '头条用户爱不爱看？',
         '9-10: 社会民生/国际政治  |  5-6: 科技需降维  |  1-2: 饭圈/小众'),
        ('天然分裂度', ACCENT_RED, 10, '天然存在两个对立群体吗？',
         '9-10: 势均力敌  |  5-6: 需转化  |  1-2: 纯感动/猎奇'),
    ]
    for name, color, weight, question, detail in dims:
        add_paragraph_rich(tf, [(name, 22, True, color), (f'  ({weight}%) -- {question}', 16, False, GRAY)], space_before=Pt(12))
        add_paragraph_rich(tf, [(detail, 14, False, GRAY)])
    add_paragraph_rich(tf, [('评级: S级(>=8分) -> 优先写  |  A级(6-8分) -> 可选  |  淘汰(<6分)', 14, False, GRAY)], space_before=Pt(16))

    # ── Slides 4-6: Topic blocks ──
    if len(topics) > 0:
        slide4 = prs.slides.add_slide(prs.slide_layouts[6])
        set_bg(slide4)
        add_textbox(slide4, 0.8, 0.2, 11, 0.5, '精选话题 (1/3)', font_size=28, bold=True, color=ACCENT_BLUE)
        draw_topic_block(slide4, topics[0], 0.5, 0.8)
        if len(topics) > 1:
            draw_topic_block(slide4, topics[1], 6.8, 0.8)

    if len(topics) > 2:
        slide5 = prs.slides.add_slide(prs.slide_layouts[6])
        set_bg(slide5)
        add_textbox(slide5, 0.8, 0.2, 11, 0.5, '精选话题 (2/3)', font_size=28, bold=True, color=ACCENT_BLUE)
        draw_topic_block(slide5, topics[2], 0.5, 0.8)
        if len(topics) > 3:
            draw_topic_block(slide5, topics[3], 6.8, 0.8)

    if len(topics) > 4:
        slide6 = prs.slides.add_slide(prs.slide_layouts[6])
        set_bg(slide6)
        add_textbox(slide6, 0.8, 0.2, 11, 0.5, '精选话题 (3/3)', font_size=28, bold=True, color=ACCENT_BLUE)
        draw_topic_block(slide6, topics[4], 0.5, 0.8)

    # ── Slide 7: Ranking ──
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_textbox(slide, 0.8, 0.2, 11, 0.5, '推荐排序', font_size=32, bold=True, color=ACCENT_BLUE)
    tf = add_rich_textbox(slide, 0.8, 0.9, 11, 5.5)
    p0 = tf.paragraphs[0]
    p0.text = ''
    medals = ['', '', '', '', '']
    
    for i, t in enumerate(topics[:5]):
        g_color = ACCENT_GREEN if t['grade'] == 'S' else ACCENT_YELLOW
        # Use custom reason if provided, otherwise generate from angle
        if 'reason' in t and t['reason']:
            reason = t['reason']
        else:
            # Generate a default reason from the angle
            angle = t.get('angle', '')
            # Truncate to ~40 chars for readability
            reason = (angle[:60] + '..') if len(angle) > 60 else angle
        
        add_paragraph_rich(tf, [
            (medals[i], 18, False, WHITE),
            (t['title'], 18, True, WHITE),
            (f'  {t["score"]}/10 ', 16, True, g_color),
            (f'{t["grade"]}级', 16, True, g_color),
        ], space_before=Pt(14))
        add_paragraph_rich(tf, [(f'      >> {reason}', 14, False, GRAY)])

    # ── Slide 8: Next Steps ──
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_textbox(slide, 0.8, 0.3, 11, 0.6, '下一步', font_size=32, bold=True, color=ACCENT_BLUE)
    tf = add_rich_textbox(slide, 0.8, 1.2, 11, 5)
    p0 = tf.paragraphs[0]
    p0.text = ''
    add_paragraph_rich(tf, [('行动建议', 24, True, ACCENT_GREEN)], space_before=Pt(10))
    add_paragraph_rich(tf, [('1. 从以上5个话题中挑选1-2个感兴趣的话题', 18, False, WHITE)], space_before=Pt(16))
    add_paragraph_rich(tf, [('2. 丢给热点刀锋进行改写/创作', 18, False, WHITE)], space_before=Pt(8))
    add_paragraph_rich(tf, [('   -> 只需说: "写一篇关于XXX的微头条"', 16, False, GRAY)], space_before=Pt(4))
    add_paragraph_rich(tf, [('3. 改写模式已被验证: 展现量是原创的20倍以上', 18, False, ACCENT_YELLOW)], space_before=Pt(16))
    add_paragraph_rich(tf, [('   (马斯克篇45万展现 vs 原创1-2万)', 16, False, GRAY)], space_before=Pt(4))
    if topics:
        add_paragraph_rich(tf, [
            ('4. 今日推荐优先选: ', 18, False, WHITE),
            (topics[0]['title'], 18, True, ACCENT_ORANGE),
            (f' ({topics[0]["score"]}/10 {topics[0]["grade"]}级)', 16, False, GRAY),
        ], space_before=Pt(16))

    return prs


def main():
    parser = argparse.ArgumentParser(description='Generate hotspot-blade topic report PPT')
    parser.add_argument('--data', help='Path to hotlist_data.json')
    parser.add_argument('--history', help='Path to hotspot-blade-history.json (deprecated, ignored)')
    parser.add_argument('--output', help='Output PPTX path (default: Desktop/选题报告_{date}.pptx)')
    parser.add_argument('--topics', help='Path to topics JSON (pre-scored topics array)')
    parser.add_argument('--today', default=datetime.now().strftime('%Y-%m-%d'),
                        help='Date string for the report (default: today)')
    args = parser.parse_args()

    # Platform counts (can be overridden by --data)
    platform_counts = {}
    if args.data:
        with open(args.data) as f:
            raw = f.read()
        idx = raw.find('{')
        data = json.loads(raw[idx:])
        for name in data:
            platform_counts[name] = data[name]['count']

    if args.topics:
        with open(args.topics) as f:
            topics = json.load(f)
    else:
        topics = json.loads(sys.stdin.read())

    today = args.today
    save_path = args.output or f'/mnt/c/Users/yingm/OneDrive/Desktop/选题报告_{today}.pptx'

    prs = build_report(topics, today=today, platform_counts=platform_counts)
    prs.save(save_path)
    print(f'PPT saved to: {save_path}')


if __name__ == '__main__':
    main()
